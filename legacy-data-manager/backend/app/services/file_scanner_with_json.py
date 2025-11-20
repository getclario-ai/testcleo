import os
import io
import re
import json
import hashlib
from datetime import datetime
from PIL import Image
import pytesseract
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from pdfminer.high_level import extract_text_to_fp
from .google_drive import GoogleDriveService
import asyncio
import logging

# Set up logging (level is configured in main.py)
logger = logging.getLogger(__name__)

# =============================================================================
# RISK SCORING SYSTEM CONFIGURATION
# =============================================================================

# Content-based risk weights for different sensitive information categories
# These weights represent the base risk level for each type of sensitive content
CONTENT_RISK_WEIGHTS = {
    'confidential': 0.4,  # Highest risk - proprietary/classified information
    'pii': 0.3,          # High risk - personally identifiable information (GDPR/HIPAA)
    'financial': 0.2,    # Medium-high risk - financial data and transactions
    'legal': 0.1          # Medium risk - legal documents and compliance
}

# Age-based risk factors (based on file creation date)
# Older files are considered riskier due to potential stale data and forgotten content
AGE_RISK_FACTORS = {
    'high': 0.3,    # Files older than 3 years - high risk (stale data)
    'medium': 0.2,  # Files 1-3 years old - medium risk
    'low': 0.1      # Files less than 1 year old - low risk
}

# Access-based risk factors (based on last accessed date)
# Files that haven't been accessed recently are considered riskier (forgotten files)
ACCESS_RISK_FACTORS = {
    'high': 0.3,    # Not accessed for >3 years - high risk (forgotten files)
    'medium': 0.2,  # Not accessed for 1-3 years - medium risk
    'low': 0.1      # Not accessed for <1 year - low risk
}

# =============================================================================
# RISK CALCULATION FUNCTIONS
# =============================================================================

def get_age_risk_factor(file_created_date):
    """
    Calculate age-based risk factor based on file creation date.
    
    Args:
        file_created_date (datetime): When the file was created
        
    Returns:
        float: Age risk factor (0.1-0.3)
    """
    if not file_created_date:
        # If no creation date available, assume medium risk
        return AGE_RISK_FACTORS['medium']
    
    try:
        # Calculate age in years
        age_years = (datetime.now() - file_created_date).days / 365
        
        if age_years >= 3:
            return AGE_RISK_FACTORS['high']    # >3 years old
        elif age_years >= 1:
            return AGE_RISK_FACTORS['medium']  # 1-3 years old
        else:
            return AGE_RISK_FACTORS['low']     # <1 year old
    except Exception as e:
        logger.warning(f"Error calculating age risk factor: {e}")
        return AGE_RISK_FACTORS['medium']

def get_access_risk_factor(last_accessed_date):
    """
    Calculate access-based risk factor based on last accessed date.
    
    Args:
        last_accessed_date (datetime): When the file was last accessed
        
    Returns:
        float: Access risk factor (0.1-0.3)
    """
    if not last_accessed_date:
        # If no access date available, assume high risk (forgotten files)
        return ACCESS_RISK_FACTORS['high']
    
    try:
        # Calculate time since last access in years
        years_since_access = (datetime.now() - last_accessed_date).days / 365
        
        if years_since_access >= 3:
            return ACCESS_RISK_FACTORS['high']    # >3 years since access
        elif years_since_access >= 1:
            return ACCESS_RISK_FACTORS['medium']  # 1-3 years since access
        else:
            return ACCESS_RISK_FACTORS['low']     # <1 year since access
    except Exception as e:
        logger.warning(f"Error calculating access risk factor: {e}")
        return ACCESS_RISK_FACTORS['medium']

def calculate_weighted_risk_score(file_data, findings):
    """
    Calculate comprehensive risk score using weighted scoring system.
    
    This function implements a multi-factor risk assessment that considers:
    1. Content-based risk (type and amount of sensitive information)
    2. Age-based risk (file creation date)
    3. Access-based risk (last accessed date)
    
    Args:
        file_data (dict): File metadata including creation and access dates
        findings (dict): Sensitive content findings by category
        
    Returns:
        float: Risk score between 0.0 and 1.0
    """
    base_score = 0.0
    
    # Step 1: Calculate content-based risk score
    # Sum up risk weights for all sensitive categories found in the file
    for category in findings.keys():
        base_score += CONTENT_RISK_WEIGHTS.get(category, 0.05)
    
    # Step 2: Add age-based risk factor
    # Older files are riskier due to potential stale data
    age_factor = get_age_risk_factor(file_data.get('createdTime'))
    base_score += age_factor
    
    # Step 3: Add access-based risk factor
    # Files not accessed recently are riskier (potentially forgotten)
    access_factor = get_access_risk_factor(file_data.get('lastAccessedTime')) # this is not available in the file metadata for g-drive file apis
    base_score += access_factor
    
    # Step 4: Cap the final score at 1.0 (100% risk)
    final_score = min(base_score, 1.0)
    
    # Log the risk calculation for debugging
    logger.debug(f"Risk calculation for {file_data.get('name', 'unknown')}: "
                f"content={base_score - age_factor - access_factor:.2f}, "
                f"age={age_factor:.2f}, access={access_factor:.2f}, "
                f"total={final_score:.2f}")
    
    return final_score

def get_risk_level_label(risk_score):
    """
    Convert numeric risk score to risk level label.
    
    Args:
        risk_score (float): Risk score between 0.0 and 1.0
        
    Returns:
        str: Risk level label ('high', 'medium', 'low')
    """
    if risk_score >= 0.7:
        return "high"
    elif risk_score >= 0.4:
        return "medium"
    else:
        return "low"

# Optional: Google API modules
try:
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload
    HAS_GOOGLE_API = True
except ImportError:
    HAS_GOOGLE_API = False

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

file_type_map = {
    'documents': ['docx', 'txt', 'doc', 'rtf', 'odt', 'pages', 'md', 'gdoc'],
    'spreadsheets': ['xlsx', 'xls', 'csv', 'ods', 'numbers', 'gsheet'],
    'presentations': ['pptx', 'ppt', 'odp', 'key', 'gslides'],
    'pdfs': ['pdf'],
    'images': ['jpg', 'jpeg', 'png', 'webp', 'gif', 'bmp', 'tiff', 'heic', 'gdraw'],
    'videos': ['mp4', 'mov', 'avi', 'wmv', 'flv', 'mkv', 'webm'],
    'audio': ['mp3', 'wav', 'ogg', 'm4a', 'wma'],
    'archives': ['zip', 'rar', '7z', 'tar', 'gz'],
    'code': ['py', 'js', 'java', 'cpp', 'h', 'cs', 'php', 'rb', 'swift', 'gs']
}

mime_type_map = {
    # Google Workspace types
    'application/vnd.google-apps.document': 'gdoc',
    'application/vnd.google-apps.spreadsheet': 'gsheet',
    'application/vnd.google-apps.presentation': 'gslides',
    'application/vnd.google-apps.drawing': 'gdraw',
    'application/vnd.google-apps.form': 'gform',
    'application/vnd.google-apps.script': 'gs',
    'application/vnd.google-apps.folder': 'folder',
    
    # Common document types
    'application/pdf': 'pdf',
    'application/msword': 'doc',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.ms-excel': 'xls',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'application/vnd.ms-powerpoint': 'ppt',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.oasis.opendocument.text': 'odt',
    'application/vnd.oasis.opendocument.spreadsheet': 'ods',
    'application/vnd.oasis.opendocument.presentation': 'odp',
    'application/x-iwork-pages-sffpages': 'pages',
    'application/x-iwork-numbers-sffnumbers': 'numbers',
    'application/x-iwork-keynote-sffkey': 'key',
    'text/markdown': 'md',
    'text/plain': 'txt',
    'text/rtf': 'rtf',
    
    # Image types
    'image/jpeg': 'jpg',
    'image/png': 'png',
    'image/gif': 'gif',
    'image/webp': 'webp',
    'image/bmp': 'bmp',
    'image/tiff': 'tiff',
    'image/heic': 'heic',
    
    # Video types
    'video/mp4': 'mp4',
    'video/quicktime': 'mov',
    'video/x-msvideo': 'avi',
    'video/x-ms-wmv': 'wmv',
    'video/webm': 'webm',
    'video/x-matroska': 'mkv',
    
    # Audio types
    'audio/mpeg': 'mp3',
    'audio/wav': 'wav',
    'audio/ogg': 'ogg',
    'audio/mp4': 'm4a',
    'audio/x-ms-wma': 'wma',
    
    # Archive types
    'application/zip': 'zip',
    'application/x-rar-compressed': 'rar',
    'application/x-7z-compressed': '7z',
    'application/x-tar': 'tar',
    'application/gzip': 'gz',
    
    # Text and code
    'text/javascript': 'js',
    'text/x-python': 'py',
    'text/x-java': 'java',
    'text/x-c': 'c',
    'text/x-cpp': 'cpp',
    'text/x-csharp': 'cs',
    'text/x-php': 'php',
    'text/x-ruby': 'rb',
    'text/x-swift': 'swift'
}

sensitive_keywords = {
    "pii": [
        "dob", "email", "phone", "address", "ssn", "personal", "pii", 
        "hipaa", "gdpr", "personally identifiable", "customer data",
        "personnel", "employee", "patient", "healthcare"
    ],
    "financial": [
        "credit", "bank", "amount", "revenue", "budget", "roi", "cost",
        "financial", "invoice", "payment", "expense", "profit", "pricing",
        "salary", "investment", "tax"
    ],
    "legal": [
        "license", "contract", "agreement", "legal", "compliance",
        "regulatory", "counsel", "policy", "policies", "terms",
        "regulation", "gdpr", "ccpa", "hipaa", "certification",
        "audit", "liability"
    ],
    "confidential": [
        "confidential", "internal use", "do not distribute", "sensitive",
        "security", "restricted", "proprietary", "classified", "private",
        "secret", "nda", "non-disclosure", "intellectual property",
        "trade secret", "internal only"
    ]
}

patterns = {
    # Matches common credit card formats (Visa, MC, Amex, Discover)
    "credit_card": r"(?:(?:4[0-9]{12}(?:[0-9]{3})?)|(?:5[1-5][0-9]{14})|(?:3[47][0-9]{13})|(?:6(?:011|5[0-9]{2})[0-9]{12}))",
    
    # Matches MM/YY or MM/YYYY with validation
    "expiry_date": r"(?:0[1-9]|1[0-2])\/(?:2[3-9]|[3-9][0-9])",
    
    # Matches SSN with required dashes and surrounding context
    "ssn": r"(?:SSN|Social Security)(?:[^0-9-])*\d{3}-\d{2}-\d{4}",
    
    # Matches email with common domains and validation
    "email": r"(?:[a-zA-Z0-9._%+-]+@(?:[a-zA-Z0-9-]+\.)+[a-zA-Z]{2,})",
    
    # Matches phone with required context and common formats
    "phone": r"(?<![\w/.:-])(?:(?:Phone|Tel|Mobile|Contact|Call|Fax)(?:[^0-9(])+)?(?:\+?1[-. ])?\(?[2-9][0-9]{2}\)?[-. ]?[2-9][0-9]{2}[-. ]?[0-9]{4}(?:\s*(?:ext|x)\.?\s*\d{1,5})?(?![-\d./@])",

    # Matches PA driver's license with validation
    "drivers_license": r"(?:Driver'?s? License|DL|License Number|License #)(?:[^0-9])*(?:[A-Z][0-9]{7}|[A-Z][0-9]{8}|[A-Z][0-9]{12}|\d{7,9}|[A-Z]\d{2}[-\s]?\d{3}[-\s]?\d{3}|[A-Z]\d{3}[-\s]?\d{3}[-\s]?\d{3}|[A-Z]{1,2}\d{4,7})",    
    
    # Matches address with validation and context
    "address_like": r"(?:Address|Location|Street)(?:[^0-9])*\d{1,5}\s[\w\s.]+(?:Street|St|Avenue|Ave|Road|Rd|Boulevard|Blvd|Lane|Ln|Drive|Dr|Circle|Cir|Court|Ct|Way|Place|Pl|Square|Sq)\b",
    
    #Legal Patterns
    "contract_number": r"(?:Contract|Agreement|License)\s*(?:#|Number|No\.?)\s*[A-Z0-9-]+",
    "legal_case": r"(?:Case|Docket)\s*(?:#|Number|No\.?)\s*[A-Z0-9-]+",
    "regulation_ref": r"(?:CFR|U\.S\.C\.|Regulation)\s+\d+[A-Z]?(?:\.\d+)*",
    
    # Confidential Patterns
    "contract_number": r"(?:Contract|Agreement|License)\s*(?:#|Number|No\.?)\s*[A-Z0-9-]+",
    "classification_level": r"(?:TOP SECRET|SECRET|CONFIDENTIAL|RESTRICTED|INTERNAL)",
    "nda_reference": r"(?:NDA|Non-Disclosure|Non Disclosure)\s+(?:Agreement|Contract)"
}

pattern_categories = {
    "credit_card": "financial",
    "expiry_date": "financial", 
    "ssn": "pii",
    "email": "pii",
    "phone": "pii",
    "drivers_license": "pii",
    "address_like": "pii",
    "contract_number": "legal",
    "legal_case": "legal", 
    "regulation_ref": "legal",
    "confidential_header": "confidential",
    "classification_level": "confidential",
    "nda_reference": "confidential"
}

now = datetime.now()

def classify_by_age(modified_time):
    age_days = (now - modified_time).days
    if age_days <= 365:
        return "lessThanOneYear"
    elif age_days <= 1095:
        return "oneToThreeYears"
    else:
        return "moreThanThreeYears"

def get_department_from_owner(owners):
    """
    Determine department based on file owner email addresses.
    
    Args:
        owners: List of owner objects from Google Drive API
        
    Returns:
        str: Department name
    """
    if not owners:
        return "Others"
    
    # Extract email addresses from owners
    owner_emails = []
    for owner in owners:
        if isinstance(owner, dict) and 'emailAddress' in owner:
            owner_emails.append(owner['emailAddress'].lower())
        elif isinstance(owner, str):
            owner_emails.append(owner.lower())
    
    # Map emails to departments
    for email in owner_emails:
        if email == "yousuf@getclario.ai":
            return "Sales & Marketing"
        elif email == "vanessa@getclario.ai":
            return "Operations"
        elif email == "madhu@getclario.ai":
            return "R&D"
    
    return "Others"

def initialize_structure():
    """Initialize the structure for file categorization."""
    return {
        "files": [],
        "stats": {
            "total_documents": 0,
            "total_sensitive": 0,
            "total_duplicates": 0,
            "by_file_type": {k: 0 for k in file_type_map.keys() | {"others"}},
            "by_sensitivity": {k: 0 for k in sensitive_keywords.keys()},
            "by_age_group": {
                "moreThanThreeYears": 0,
                "oneToThreeYears": 0,
                "lessThanOneYear": 0
            },
            "by_risk_level": {
                "high": 0,
                "medium": 0,
                "low": 0
            }
        }
    }

def scan_text(text):
    """
    Scan text for sensitive information using keywords and patterns.
    Returns a dictionary of findings only if sensitive content is detected.
    """
    findings = {cat: [] for cat in sensitive_keywords}
    text_lower = text.lower()
    
    # Check for keywords in each category
    for cat, keywords in sensitive_keywords.items():
        for keyword in keywords:
            # Look for whole word matches only
            if re.search(r'\b' + re.escape(keyword.lower()) + r'\b', text_lower):
                findings[cat].append(keyword)
    
    # Check for pattern matches
    for label, pattern in patterns.items():
        if re.search(pattern, text):
            category = pattern_categories.get(label, "pii")  # Default to pii if not specified
            findings[category].append(label)
    
    # Only return categories that have findings
    return {k: v for k, v in findings.items() if v}

def extract_text_from_file(stream, file_type): 
    try:
        if file_type == 'docx':
            doc = Document(stream)
            return "\n".join([p.text for p in doc.paragraphs])
        elif file_type == 'pptx':
            prs = Presentation(stream)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file_type in ['xlsx', 'xls']:
            wb = load_workbook(stream, read_only=True, data_only=True)
            return "\n".join([str(cell.value) for sheet in wb.worksheets for row in sheet.iter_rows() for cell in row if cell.value])
        elif file_type == 'pdf':
            output = io.StringIO()
            extract_text_to_fp(stream, output)
            return output.getvalue()
        elif file_type in ['jpg', 'jpeg', 'png', 'webp']:
            img = Image.open(stream)
            return pytesseract.image_to_string(img)
        elif file_type == 'txt':
            return stream.read().decode()
    except:
        return ""
    return ""

async def scan_files(source='local', path_or_drive_id='.', output_json='scan_report.json', drive_service=None):
    # Initialize a simplified, flattened data structure
    results = {
        "files": [],
        "stats": {
            "total_documents": 0,
            "total_sensitive": 0,
            "total_duplicates": 0,
            "by_file_type": {k: 0 for k in file_type_map.keys() | {"others"}},
            "by_sensitivity": {k: 0 for k in sensitive_keywords.keys()},
            "by_age_group": {
                "moreThanThreeYears": 0,
                "oneToThreeYears": 0,
                "lessThanOneYear": 0
            },
            "by_risk_level": {
                "high": 0,
                "medium": 0,
                "low": 0
            },
            "by_department": {}
        },
        "scan_complete": False,
        "processed_files": 0,
        "total_files": 0,
        "failed_files": []
    }

    # Add logging for file type mapping
    logger.debug(f"Using file type mapping: {file_type_map}")
    type_counts = {k: 0 for k in file_type_map.keys() | {"others"}}

    if source == 'local':
        all_files = []
        for root, _, files in os.walk(path_or_drive_id):
            for f in files:
                all_files.append(os.path.join(root, f))
        results["total_files"] = len(all_files)

        for filepath in all_files:
            try:
                ext = filepath.split('.')[-1].lower()
                modified_time = datetime.fromtimestamp(os.path.getmtime(filepath))
                age_group = classify_by_age(modified_time)
                file_type = next((k for k, v in file_type_map.items() if ext in v), "others")
                with open(filepath, 'rb') as f:
                    content = extract_text_from_file(f, ext)
                if not content:
                    continue
                results[age_group]["total_documents"] += 1
                results[age_group]["file_types"][file_type].append(filepath)
                findings = scan_text(content)
                if findings:
                    results[age_group]["total_sensitive"] += 1
                    results["total_sensitive_files"] += 1
                    for k, v in findings.items():
                        results[age_group]["sensitive_info"][k].extend(v)
                results["processed_files"] += 1
            except Exception as e:
                logger.error(f"Error processing file {filepath}: {str(e)}")
                results["failed_files"].append(filepath)

    elif source == 'gdrive' and HAS_GOOGLE_API:
        # Use provided drive_service if available, otherwise create a new one (for backward compatibility)
        if drive_service is None:
            drive_service = GoogleDriveService()
        if not await drive_service.is_authenticated():
            raise ValueError("Not authenticated with Google Drive")

        try:
            files = await drive_service.list_directory(path_or_drive_id, recursive=True)
            results["total_files"] = len(files)
            
            # Track unique sensitive files
            sensitive_file_ids = set()

            for file in files:
                try:
                    file_id = file['id']
                    name = file['name']
                    mime_type = file['mimeType']
                    
                    # Log file type categorization
                    logger.debug(f"Processing file: {name} (mime_type: {mime_type})")
                    
                    # Get file extension from mime type or name
                    ext = mime_type_map.get(mime_type, None)
                    if not ext and '.' in name:
                        ext = name.split('.')[-1].lower()
                    
                    if not ext:
                        ext = 'others'

                    modified_time = datetime.fromisoformat(file['modifiedTime'].rstrip("Z"))
                    age_group = classify_by_age(modified_time)

                    # Determine file type category
                    file_type = 'others'
                    for category, extensions in file_type_map.items():
                        if ext in extensions:
                            file_type = category
                            break
                    
                    # Update type counts
                    type_counts[file_type] += 1
                    
                    # Determine department from file owners
                    department = get_department_from_owner(file.get('owners', []))
                    
                    # Create a standardized file object
                    file_dict = {
                        "id": file_id,
                        "name": name,
                        "mimeType": mime_type,
                        "modifiedTime": file['modifiedTime'],
                        "createdTime": file.get('createdTime', file['modifiedTime']),
                        "size": int(file.get('size', 0)),
                        "fileType": file_type,
                        "ageGroup": age_group,
                        "sensitiveCategories": [],
                        "riskLevel": None,
                        "riskLevelLabel": None,
                        "department": department
                    }
                    
                    # Add to the flat files array
                    results["files"].append(file_dict)
                    
                    # Update statistics
                    results["stats"]["by_file_type"][file_type] += 1
                    results["stats"]["by_age_group"][age_group] += 1
                    
                    # Update department statistics
                    if department not in results["stats"]["by_department"]:
                        results["stats"]["by_department"][department] = 0
                    results["stats"]["by_department"][department] += 1

                    # Only scan content for text-based files
                    if file_type in ['documents', 'spreadsheets', 'presentations', 'pdfs']:
                        try:
                            content = await drive_service.get_file_content(file_id)
                            if content:
                                findings = scan_text(content)
                                if findings:  # If any sensitive content was found
                                    # Find the file in our results array
                                    file_index = next((i for i, f in enumerate(results["files"]) if f["id"] == file_id), None)
                                    if file_index is not None:
                                        # Update the file with sensitivity information
                                        file_obj = results["files"][file_index]
                                        
                                        # Track all sensitive categories found
                                        all_categories = []
                                        explanations = []
                                        
                                        for category, found_items in findings.items():
                                            if found_items:  # Only add if there are findings
                                                all_categories.append(category)
                                                explanations.append(f"Found {', '.join(found_items)} in {category} category")
                                                # Update stats
                                                results['stats']['by_sensitivity'][category] += 1
                                        
                                        # Update the file object with sensitivity data
                                        file_obj["sensitiveCategories"] = all_categories
                                        file_obj["sensitivityExplanation"] = "; ".join(explanations)
                                        file_obj["confidence"] = 0.8
                                        
                                        # Mark as sensitive in stats if not already counted
                                        if file_id not in sensitive_file_ids:
                                            results["stats"]["total_sensitive"] += 1
                                            sensitive_file_ids.add(file_id)
                        except Exception as e:
                            logger.error(f"Error processing file content {name}: {str(e)}")
                    
                    results["processed_files"] += 1
                except Exception as e:
                    logger.error(f"Error processing file {name}: {str(e)}")
                    results["failed_files"].append(name)

            logger.debug(f"Completed processing {results['processed_files']} files")
            logger.debug(f"Found {len(sensitive_file_ids)} sensitive files")
            results["scan_complete"] = True
            
            # --- CALCULATE WEIGHTED RISK SCORES FOR SENSITIVE FILES ---
            # Now calculate risk scores for all files with sensitive content
            for file_obj in results["files"]:
                if file_obj["sensitiveCategories"]:
                    # Get file creation time and last accessed time for risk calculation
                    file_data = {
                        "createdTime": datetime.fromisoformat(file_obj["createdTime"].rstrip("Z")) if "createdTime" in file_obj else None,
                        "lastAccessedTime": None,  # We don't have this data from Google Drive API
                        "name": file_obj["name"]
                    }
                    
                    # Convert sensitiveCategories to the format expected by calculate_weighted_risk_score
                    findings_format = {category: [category] for category in file_obj["sensitiveCategories"]}
                    
                    # Calculate risk score using our weighted scoring system
                    risk_score = calculate_weighted_risk_score(file_data, findings_format)
                    risk_level_label = get_risk_level_label(risk_score)
                    
                    # Update the file object with risk data
                    file_obj["riskLevel"] = risk_score
                    file_obj["riskLevelLabel"] = risk_level_label
                    
                    # Update risk level stats
                    results["stats"]["by_risk_level"][risk_level_label] += 1
                    
                    # Log the risk assessment for debugging
                    logger.debug(f"File {file_obj.get('name', 'unknown')} risk assessment: "
                               f"score={risk_score:.2f}, level={risk_level_label}, "
                               f"categories={file_obj['sensitiveCategories']}")
            
            # --- UPDATE OVERALL STATISTICS ---
            # Now that we've processed all files, update the overall statistics
            results["stats"]["total_documents"] = len(results["files"])
            results["total_files"] = len(results["files"])
            results["processed_files"] = len(results["files"])
            results["stats"]["total_sensitive"] = len(sensitive_file_ids)
            
            # Log summary of risk distribution for the entire scan
            high_risk_files = sum(1 for file in results["files"] if file.get("riskLevelLabel") == "high")
            medium_risk_files = sum(1 for file in results["files"] if file.get("riskLevelLabel") == "medium")
            low_risk_files = sum(1 for file in results["files"] if file.get("riskLevelLabel") == "low")
            
            logger.debug(f"Risk distribution summary: High={high_risk_files}, Medium={medium_risk_files}, Low={low_risk_files}")
            logger.debug(f"Total sensitive files processed: {len(sensitive_file_ids)}")
            
            # Add primary sensitivity reason to each file (highest weighted category)
            for file_obj in results["files"]:
                if file_obj["sensitiveCategories"]:
                    # Determine primary sensitivity reason (highest weighted category)
                    primary_category = max(file_obj["sensitiveCategories"], 
                                          key=lambda cat: CONTENT_RISK_WEIGHTS.get(cat, 0), 
                                          default=None)
                    file_obj["sensitivityReason"] = primary_category
            # No backward compatibility layer needed
       
        except Exception as e:
            logger.error(f"Error scanning files: {str(e)}")
            raise

    return results
